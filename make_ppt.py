#!/usr/bin/env python3
"""
PPT 演示文稿生成脚本
用法: python make_ppt.py --spec spec.json --output presentation.pptx
功能: 读取 JSON 规格文件 → 生成 16:9 宽屏 PowerPoint 演示文稿
"""

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

try:
    from PIL import Image
except ImportError:
    Image = None

# ── 常量 ──────────────────────────────────────────────

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# 颜色
COLOR_BG_DARK = RGBColor(0x1F, 0x2B, 0x3E)      # 深蓝夜色标题页背景
COLOR_TITLE_BAR = RGBColor(0x2D, 0x5F, 0x8A)     # 中蓝标题栏
COLOR_ACCENT = RGBColor(0xE8, 0x7C, 0x2A)        # 橙色点缀
COLOR_TEXT_DARK = RGBColor(0x33, 0x33, 0x33)      # 深灰正文
COLOR_TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)     # 白色（深色背景上）
COLOR_TEXT_LIGHT = RGBColor(0xAA, 0xBB, 0xCC)     # 浅灰副标题
COLOR_OVERLAY = RGBColor(0x0A, 0x14, 0x20)        # 图片页底部遮罩

# 字体
FONT_LATIN = "Arial"
FONT_CJK_MAC = "PingFang SC"
FONT_CJK_WIN = "Microsoft YaHei"
FONT_CJK = FONT_CJK_MAC if sys.platform == "darwin" else FONT_CJK_WIN


# ── 工具函数 ──────────────────────────────────────────

def has_cjk(text: str) -> bool:
    """检测文本是否包含中文字符"""
    for ch in text:
        if "一" <= ch <= "鿿":
            return True
    return False


def set_run_font(run, text: str, size: int, bold: bool = False, color: RGBColor = COLOR_TEXT_DARK):
    """设置文本 run 的字体属性，自动处理中英文"""
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

    if has_cjk(text):
        run.font.name = FONT_CJK
        # 设置东亚文字字体
        rPr = run._r.get_or_add_rPr()
        ea = rPr.makeelement(qn("a:ea"), {})
        ea.set("typeface", FONT_CJK)
        rPr.append(ea)
    else:
        run.font.name = FONT_LATIN


def add_textbox(slide, left, top, width, height):
    """添加文本框并返回 (textbox, text_frame)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def add_title_text(text_frame, title: str, size: int, color: RGBColor = COLOR_TEXT_WHITE, bold: bool = True, align=PP_ALIGN.LEFT):
    """向 text_frame 添加标题段落"""
    p = text_frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    run = p.add_run()
    run.text = title
    set_run_font(run, title, size, bold=bold, color=color)
    return p


def add_bullet_paragraph(text_frame, bullet: str, size: int = 20, color: RGBColor = COLOR_TEXT_DARK):
    """向 text_frame 追加一个要点段落"""
    p = text_frame.add_paragraph()
    p.level = 0
    p.space_after = Pt(8)
    p.space_before = Pt(4)
    # 手动添加 bullet 字符
    run = p.add_run()
    run.text = f"•  {bullet}"
    set_run_font(run, bullet, size, bold=False, color=color)
    return p


def fill_background(slide, color: RGBColor):
    """填充幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill_color: RGBColor):
    """添加一个矩形色块"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()  # 无边框
    return shape


def add_image_scaled(slide, image_path: str, left, top, max_width, max_height):
    """
    将图片等比缩放后居中放置到指定区域。
    返回 True 如果成功，False 如果图片不存在或无法读取。
    """
    path = Path(image_path)
    if not path.exists():
        print(f"WARNING: 图片不存在: {image_path}", file=sys.stderr)
        return False

    try:
        if Image is not None:
            with Image.open(image_path) as img:
                img_width, img_height = img.size
        else:
            # 无 Pillow，直接按 max_width 放置，让 python-pptx 处理缩放
            slide.shapes.add_picture(image_path, left, top, max_width)
            return True

        aspect = img_width / img_height
        box_aspect = max_width / max_height

        if box_aspect > aspect:
            # 区域更宽，以高度为约束
            final_height = max_height
            final_width = int(max_height * aspect)
        else:
            # 区域更高，以宽度为约束
            final_width = max_width
            final_height = int(max_width / aspect)

        # 居中偏移
        final_left = left + (max_width - final_width) // 2
        final_top = top + (max_height - final_height) // 2

        slide.shapes.add_picture(
            image_path, final_left, final_top, final_width, final_height
        )
        return True

    except Exception as e:
        print(f"WARNING: 无法添加图片 {image_path}: {e}", file=sys.stderr)
        return False


# ── 幻灯片创建函数 ───────────────────────────────────

def create_title_slide(prs: Presentation, spec: dict):
    """标题页：深色背景 + 居中标题 + 副标题 + 装饰线"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    fill_background(slide, COLOR_BG_DARK)

    # 装饰条 - 顶部
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), COLOR_ACCENT)

    # 主标题
    title = spec.get("title", "")
    _, tf = add_textbox(slide, Inches(1.5), Inches(2.2), Inches(10.3), Inches(2.0))
    add_title_text(tf, title, size=44, color=COLOR_TEXT_WHITE, align=PP_ALIGN.CENTER)

    # 副标题
    subtitle = spec.get("subtitle", "")
    if subtitle:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(16)
        run = p.add_run()
        run.text = subtitle
        set_run_font(run, subtitle, size=22, bold=False, color=COLOR_TEXT_LIGHT)

    # 装饰线
    add_rect(slide, Inches(5.5), Inches(4.5), Inches(2.3), Inches(0.05), COLOR_ACCENT)

    # 作者信息
    author = spec.get("author", "")
    if author:
        _, tf2 = add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0))
        p = tf2.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = author
        set_run_font(run, author, size=16, bold=False, color=COLOR_TEXT_LIGHT)


def create_content_slide(prs: Presentation, spec: dict):
    """内容页：标题栏 + 要点列表"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题栏背景
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), COLOR_TITLE_BAR)

    # 标题文字
    title = spec.get("title", "")
    _, tf = add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.9))
    add_title_text(tf, title, size=30, color=COLOR_TEXT_WHITE)

    # 要点列表
    bullets = spec.get("bullets", [])
    if bullets:
        _, tf2 = add_textbox(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(5.0))
        tf2.paragraphs[0].text = ""  # 清空默认段落
        for i, bullet in enumerate(bullets):
            add_bullet_paragraph(tf2, bullet, size=20, color=COLOR_TEXT_DARK)

    # 底部装饰线
    add_rect(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.04), COLOR_ACCENT)


def create_content_image_slide(prs: Presentation, spec: dict):
    """图文页：标题栏 + 要点 + 配图（image_position 控制图片在左/右）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # 标题栏
    add_rect(slide, Inches(0), Inches(0), SLIDE_WIDTH, Inches(1.2), COLOR_TITLE_BAR)

    title = spec.get("title", "")
    _, tf = add_textbox(slide, Inches(0.8), Inches(0.2), Inches(11.7), Inches(0.9))
    add_title_text(tf, title, size=30, color=COLOR_TEXT_WHITE)

    # 图文布局
    image_pos = spec.get("image_position", "right")
    bullets = spec.get("bullets", [])
    image_path = spec.get("image", "")

    # 文本区域宽度
    text_width = Inches(5.8)
    image_width = Inches(5.5)
    margin = Inches(0.8)
    gap = Inches(0.4)

    if image_pos == "left":
        text_left = margin + image_width + gap
        image_left = margin
    else:  # right (默认)
        text_left = margin
        image_left = margin + text_width + gap

    # 要点
    if bullets:
        _, tf2 = add_textbox(slide, text_left, Inches(1.8), text_width, Inches(5.0))
        tf2.paragraphs[0].text = ""
        for bullet in bullets:
            add_bullet_paragraph(tf2, bullet, size=18, color=COLOR_TEXT_DARK)

    # 图片
    if image_path:
        add_image_scaled(
            slide, image_path,
            image_left, Inches(1.8),
            image_width, Inches(5.0)
        )

    # 底部装饰线
    add_rect(slide, Inches(0), Inches(7.35), SLIDE_WIDTH, Inches(0.04), COLOR_ACCENT)


def create_image_slide(prs: Presentation, spec: dict):
    """图片页：全幅图片 + 底部半透明说明条"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_background(slide, COLOR_BG_DARK)

    image_path = spec.get("image", "")
    caption = spec.get("caption", "")

    # 全幅图片
    if image_path:
        add_image_scaled(
            slide, image_path,
            Inches(0), Inches(0),
            SLIDE_WIDTH, SLIDE_HEIGHT
        )

    # 底部遮罩（半透明深色）
    overlay = add_rect(slide, Inches(0), Inches(5.5), SLIDE_WIDTH, Inches(2.0), COLOR_OVERLAY)
    # 设置半透明（alpha 70%）通过 XML 操作
    spPr = overlay._element.find(qn("p:spPr"))
    if spPr is not None:
        solidFill = spPr.find(qn("a:solidFill"))
        if solidFill is not None:
            srgbClr = solidFill.find(qn("a:srgbClr"))
            if srgbClr is not None:
                alpha_elem = srgbClr.makeelement(qn("a:alpha"), {"val": "70000"})
                srgbClr.append(alpha_elem)

    # 说明文字
    if caption:
        _, tf = add_textbox(slide, Inches(1.0), Inches(5.8), Inches(11.3), Inches(1.2))
        add_title_text(tf, caption, size=24, color=COLOR_TEXT_WHITE, align=PP_ALIGN.LEFT)


# ── 规格验证 ──────────────────────────────────────────

VALID_TYPES = {"title", "content", "content_image", "image"}


def validate_spec(spec: dict) -> bool:
    """验证规格文件，打印错误信息，返回是否有效"""
    valid = True

    if "slides" not in spec:
        print("ERROR: 规格文件缺少 'slides' 字段", file=sys.stderr)
        return False

    for i, slide in enumerate(spec["slides"]):
        slide_type = slide.get("type")
        if slide_type not in VALID_TYPES:
            print(f"ERROR: 幻灯片 {i} 的类型 '{slide_type}' 无效，应为 {VALID_TYPES}", file=sys.stderr)
            valid = False

        if slide_type == "content_image" and "image" not in slide:
            print(f"ERROR: 幻灯片 {i} 是 content_image 类型但缺少 'image' 路径", file=sys.stderr)
            valid = False

        if slide_type == "image" and "image" not in slide:
            print(f"ERROR: 幻灯片 {i} 是 image 类型但缺少 'image' 路径", file=sys.stderr)
            valid = False

        # 图片存在性检查（WARNING 级别，不阻断）
        if slide_type in ("content_image", "image"):
            img = slide.get("image", "")
            if img and not Path(img).exists():
                print(f"WARNING: 幻灯片 {i} 的图片不存在: {img}", file=sys.stderr)

    return valid


# ── 主流程 ────────────────────────────────────────────

SLIDE_BUILDERS = {
    "title": create_title_slide,
    "content": create_content_slide,
    "content_image": create_content_image_slide,
    "image": create_image_slide,
}


def main():
    parser = argparse.ArgumentParser(description="PPT 演示文稿生成脚本")
    parser.add_argument("--spec", required=True, help="JSON 规格文件路径")
    parser.add_argument("--output", required=True, help="输出 .pptx 文件路径")
    args = parser.parse_args()

    # 读取规格
    spec_path = Path(args.spec)
    if not spec_path.exists():
        print(f"ERROR: 规格文件不存在: {args.spec}", file=sys.stderr)
        sys.exit(1)

    with open(spec_path, "r", encoding="utf-8") as f:
        spec = json.load(f)

    # 验证
    if not validate_spec(spec):
        sys.exit(1)

    # 创建演示文稿
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # 按规格逐页创建
    for i, slide_spec in enumerate(spec.get("slides", [])):
        slide_type = slide_spec.get("type", "content")
        builder = SLIDE_BUILDERS.get(slide_type)
        if builder:
            builder(prs, slide_spec)
            print(f"  ✓ 幻灯片 {i + 1}: {slide_type}", file=sys.stderr)
        else:
            print(f"  ✗ 幻灯片 {i + 1}: 未知类型 '{slide_type}'，跳过", file=sys.stderr)

    # 保存
    output_path = Path(args.output)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))

    # 打印输出路径（供 OpenClaw 读取）
    print(str(output_path.resolve()))


if __name__ == "__main__":
    main()
